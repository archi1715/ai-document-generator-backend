from fastapi import APIRouter, HTTPException, Depends, Query
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
from bson import ObjectId
from fastapi.encoders import jsonable_encoder
from app.auth.dependencies import get_current_user
from app.db.mongo import get_documents_collection, get_document_history_collection
from app.models.document import DocumentHistory

def generate_ppt_doc(content: str) -> BytesIO:
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = "AI Generated Content"
    body.text = content

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

router = APIRouter(prefix="/api/ppt", tags=["PPT Generator"])

#get history document by id 
@router.get("/history/{document_id}", response_model=DocumentHistory)
async def get_document_by_id(document_id: str, user=Depends(get_current_user)):
    collection = get_document_history_collection()
    document = await collection.find_one({
        "_id": ObjectId(document_id),
        "user_email": user["email"]
    })

    if not document:
        raise HTTPException(status_code=404, detail="Document not found")

    return jsonable_encoder(document)

#delete historical document by id
@router.delete("/history/{document_id}")
async def delete_document_by_id(document_id: str, user=Depends(get_current_user)):
    collection = get_document_history_collection()
    result = await collection.delete_one({
        "_id": ObjectId(document_id),
        "user_email": user["email"]
    })

    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Document not found or not authorized")

    return {"message": "Document deleted successfully"}
